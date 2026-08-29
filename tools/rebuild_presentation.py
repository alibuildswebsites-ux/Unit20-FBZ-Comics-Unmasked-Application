from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_LINE_DASH_STYLE

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / 'submission' / 'Unit20_FBZ_SOLID_OOP_Presentation.pptx'
OLD = ROOT / 'submission' / 'Unit20_FBZ_SOLID_OOP_Presentation.pptx.bak'
OLD.write_bytes(OUT.read_bytes()) if not OLD.exists() else None

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

INK = RGBColor(25,35,45)
MUTED = RGBColor(90,100,112)
PALE = RGBColor(245,247,250)
WHITE = RGBColor(255,255,255)
ACCENT = RGBColor(60,90,120)
ACCENT2 = RGBColor(110,140,165)
GREEN = RGBColor(75,130,95)
RED = RGBColor(170,85,85)
GOLD = RGBColor(170,130,55)


def text(slide, x,y,w,h, value, size=22, bold=False, color=INK, align=PP_ALIGN.LEFT, valign=MSO_VERTICAL_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear(); tf.word_wrap = True; tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = value; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
    r.font.name = 'Aptos'
    return tb


def box(slide,x,y,w,h,title,body='',fill=PALE,line=ACCENT, title_size=22, body_size=15):
    sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x),Inches(y),Inches(w),Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill; sh.line.color.rgb=line; sh.line.width=Pt(1.5)
    text(slide,x+0.15,y+0.12,w-0.3,0.42,title,title_size,True,INK)
    if body: text(slide,x+0.15,y+0.62,w-0.3,h-0.72,body,body_size,False,MUTED)
    return sh


def line(slide,x1,y1,x2,y2,color=ACCENT,width=2,dash=None):
    c=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1),Inches(y1),Inches(x2),Inches(y2))
    c.line.color.rgb=color; c.line.width=Pt(width)
    if dash: c.line.dash_style=dash
    return c


def footer(slide, n, section='Unit 20 • FBZ'):
    text(slide,0.7,7.1,6,0.2,f'{section} • Slide {n}',9,False,MUTED)


def title(slide, kicker, head, sub=None):
    text(slide,0.7,0.3,11.9,0.3,kicker.upper(),10,True,ACCENT)
    text(slide,0.7,0.65,12,0.65,head,28,True,INK)
    if sub: text(slide,0.7,1.27,11.8,0.4,sub,13,False,MUTED)

# 1
s=prs.slides.add_slide(blank)
text(s,0.75,0.55,5.2,0.35,'UNIT 20 • APPLIED PROGRAMMING',11,True,ACCENT)
text(s,0.75,1.25,7.2,1.25,'Fantasy Bazaar\nComics Unmasked',32,True,INK)
text(s,0.78,2.75,6.8,0.6,'SOLID OOP • design patterns • real-data validation • automated testing',16,False,MUTED)
box(s,8.5,1.35,3.85,1.45,'REAL DATA','57,746 records view\n117,873 names rows',fill=PALE,line=ACCENT)
box(s,8.5,3.05,3.85,1.45,'FINAL EVIDENCE','33 / 33 tests passing\n5 official CSV views validated',fill=PALE,line=GREEN)
box(s,8.5,4.75,3.85,1.45,'DESIGN OUTCOME','Repository + Strategy + Factory\nrole-aware domain model',fill=PALE,line=GOLD)
text(s,0.78,6.45,7.2,0.6,'A requirement-led implementation—not SOLID for its own sake.',18,True,INK)
footer(s,1)

# 2
s=prs.slides.add_slide(blank); title(s,'01 • scenario','What problem is the application solving?','The project turns a large catalogue dataset into a maintainable, testable desktop/CLI application.')
box(s,0.75,2.0,3.65,2.0,'INPUT','British Library Comics Unmasked\n2022 researcher-format CSV package\nFive linked/faceted views',fill=WHITE)
box(s,4.85,2.0,3.65,2.0,'PROCESS','Load → model → aggregate → search\nGroup by author/year\nSort A–Z / Z–A',fill=WHITE)
box(s,8.95,2.0,3.65,2.0,'OUTPUT','User-facing catalogue records\nAdvanced search\nReports + acceptance evidence',fill=WHITE)
line(s,4.42,3.0,4.83,3.0); line(s,8.52,3.0,8.93,3.0)
text(s,0.8,4.55,11.5,0.65,'Core design decision: preserve source fidelity while presenting the user with one logical record per BL record ID.',19,True,INK)
footer(s,2)

# 3
s=prs.slides.add_slide(blank); title(s,'02 • dataset','The real dataset changed the architecture','The implementation was refined after inspecting the actual Comics Unmasked facet structure.')
# metric row
metrics=[('57,746','records.csv'),('117,873','names.csv'),('77,280','titles.csv'),('77,919','topics.csv'),('57,844','classification.csv')]
for i,(num,lab) in enumerate(metrics): box(s,0.7+i*2.5,2.0,2.25,1.2,lab,num,fill=PALE,line=ACCENT,title_size=11,body_size=21)
text(s,0.8,3.55,5.9,0.55,'Names view: 117,873 facet rows',23,True,INK)
text(s,0.8,4.15,5.9,0.8,'→ 54,147 unique BL record IDs\n→ 63,726 repeated/facet rows collapsed',21,True,INK)
box(s,7.1,3.45,5.0,2.1,'WHY THIS MATTERS','A flat row-by-row display would duplicate logical catalogue entries.\n\nAggregation became a first-class repository responsibility.',fill=WHITE,line=GOLD,title_size=20,body_size=16)
footer(s,3)

# 4
s=prs.slides.add_slide(blank); title(s,'03 • architecture','Layered architecture keeps change in the right place')
# layers
layers=[('Presentation','CLI / interactive workflow',0.75,1.85,11.8),('Application services','EncyclopediaService • SearchService • FavouriteService',1.15,2.65,11.0),('Domain','Comic • Contributor • SearchCriteria',1.55,3.45,10.2),('Patterns','Strategy • Simple Factory • Repository',1.95,4.25,9.4),('Data sources','CSV repositories • aggregation • XML adapter • in-memory repository',2.35,5.05,8.6)]
for label,body,x,y,w in layers: box(s,x,y,w,0.62,label,body,fill=PALE,line=ACCENT,title_size=14,body_size=12)
text(s,0.8,6.0,11.4,0.55,'Boundary rule: services depend on abstractions; the CLI does not own business logic; data-source details stay below the repository boundary.',16,True,INK)
footer(s,4)

# 5
s=prs.slides.add_slide(blank); title(s,'04 • oop model','OOP relationships in the actual code')
box(s,0.7,2.0,3.5,1.55,'Comic','Immutable catalogue record\nDomain operations: tokens(), authors()',fill=WHITE,line=ACCENT)
box(s,4.9,2.0,3.5,1.55,'Contributor','name + role\nExplicit author/writer role detection',fill=WHITE,line=ACCENT)
box(s,9.1,2.0,3.5,1.55,'ComicRepository','Abstract contract: all()\nCSV / XML / memory implementations',fill=WHITE,line=ACCENT)
box(s,2.8,4.55,3.5,1.4,'SearchStrategy','Common search(comics, query) contract',fill=WHITE,line=GOLD)
box(s,7.0,4.55,3.5,1.4,'SearchService','Depends on repository\nDelegates to strategy',fill=WHITE,line=GOLD)
line(s,4.2,2.77,4.82,2.77); line(s,8.42,2.77,9.02,2.77)
line(s,6.55,5.25,6.95,5.25); line(s,5.9,3.58,5.9,4.5); line(s,7.1,3.58,8.5,4.48)
text(s,0.8,6.35,11.5,0.4,'Most important semantic refinement: preserving contributor role prevents “editor = author” false matches.',16,True,INK)
footer(s,5)

# 6
s=prs.slides.add_slide(blank); title(s,'05 • solid','SOLID is visible in concrete change boundaries')
solid=[('S','Single Responsibility','CSV parsing, aggregation, search, favourites, presentation'),('O','Open / Closed','Add a search strategy without rewriting the service'),('L','Liskov Substitution','CSV, XML, memory repositories honour the same contract'),('I','Interface Segregation','Repository and Strategy interfaces stay small'),('D','Dependency Inversion','Services depend on ComicRepository, not CSV implementation')]
for i,(k,h,b) in enumerate(solid):
    y=1.95+i*0.9; box(s,0.75,y,0.75,0.72,k,fill=PALE,line=ACCENT,title_size=25); text(s,1.7,y,2.8,0.3,h,15,True,INK); text(s,4.1,y,8.3,0.52,b,14,False,MUTED)
text(s,0.8,6.6,11.5,0.35,'Design principle: each abstraction exists because the application has a real variation or testing need.',15,True,INK)
footer(s,6)

# 7
s=prs.slides.add_slide(blank); title(s,'06 • design patterns','Three patterns/techniques solve three different problems')
box(s,0.75,2.0,3.65,2.2,'STRATEGY','Behavioral variation\n\nTitle • Author • Genre • Year\nSame search contract, different algorithms',fill=WHITE,line=ACCENT,title_size=20,body_size=15)
box(s,4.85,2.0,3.65,2.2,'SIMPLE FACTORY','Creation variation\n\nSearchStrategyFactory maps a search type to a concrete strategy.',fill=WHITE,line=GOLD,title_size=20,body_size=15)
box(s,8.95,2.0,3.65,2.2,'REPOSITORY','Storage variation\n\nCSV, XML and memory sources sit behind ComicRepository.',fill=WHITE,line=GREEN,title_size=20,body_size=15)
text(s,0.8,4.8,11.5,0.8,'Technical accuracy matters: the project calls the creation technique a “Simple Factory,” not the formal GoF Factory Method.',17,True,INK)
footer(s,7)

# 8
s=prs.slides.add_slide(blank); title(s,'07 • pipeline','From raw CSV to user-visible result')
steps=[('1','CSV row','raw strings'),('2','Comic','validated domain object'),('3','Aggregate','one BL ID record'),('4','Search','strategy/service filters'),('5','Sort','A–Z or Z–A'),('6','Present','CLI / report')]
for i,(n,h,b) in enumerate(steps):
    x=0.55+i*2.1; box(s,x,2.25,1.8,1.5,h,b,fill=PALE,line=ACCENT,title_size=17,body_size=12); text(s,x+0.65,1.7,0.5,0.35,n,15,True,ACCENT,PP_ALIGN.CENTER)
    if i<5: line(s,x+1.8,3.0,x+2.05,3.0)
text(s,0.8,4.45,11.3,0.75,'Validation happens at the domain boundary; business logic operates on stable objects; presentation receives prepared results.',19,True,INK)
box(s,0.8,5.55,11.3,0.85,'DATA SAFETY','Identifiers remain strings (including leading zeros); missing values become “missing”; multi-value fields remain semantically separate.',fill=WHITE,line=GOLD,title_size=14,body_size=14)
footer(s,8)

# 9
s=prs.slides.add_slide(blank); title(s,'08 • search','Search is composable instead of conditional')
box(s,0.75,2.0,2.5,1.5,'SearchService','Coordinates repository + strategy',fill=WHITE,line=ACCENT,title_size=18,body_size=14)
box(s,4.0,1.65,2.2,0.95,'Title','substring title match',fill=PALE,line=ACCENT,title_size=15,body_size=11)
box(s,4.0,2.8,2.2,0.95,'Author','explicit author/writer roles',fill=PALE,line=ACCENT,title_size=15,body_size=11)
box(s,4.0,3.95,2.2,0.95,'Genre','token-aware genre match',fill=PALE,line=ACCENT,title_size=15,body_size=11)
box(s,6.9,2.2,2.2,0.95,'Year','publication-year search',fill=PALE,line=ACCENT,title_size=15,body_size=11)
box(s,9.8,2.2,2.65,2.4,'Future extension','ISBN, publisher, language, etc. can be added as new strategies or focused service filters.',fill=WHITE,line=GREEN,title_size=18,body_size=14)
line(s,3.25,2.75,4.0,2.15); line(s,3.25,2.75,4.0,3.25); line(s,3.25,2.75,4.0,4.35); line(s,3.25,2.75,6.9,2.65); line(s,9.1,2.65,9.8,3.1)
footer(s,9)

# 10
s=prs.slides.add_slide(blank); title(s,'09 • aggregation','Why aggregation is not just a cosmetic step')
text(s,0.85,1.95,5.0,0.45,'117,873 raw names rows',25,True,INK)
text(s,0.85,2.55,5.0,0.45,'− 63,726 repeated/facet rows',22,True,MUTED)
text(s,0.85,3.1,5.0,0.45,'= 54,147 logical catalogue records',25,True,GREEN)
box(s,6.7,2.0,5.3,2.8,'AggregationRepository','Groups by BL record ID\n\nMerges repeated metadata values\nPreserves variant titles\nPreserves Contributor(name, role) pairs\nLeaves raw source files untouched',fill=WHITE,line=GOLD,title_size=20,body_size=15)
text(s,0.85,4.4,11.2,0.7,'The design follows the data model: the repository adapts a facet-oriented source into a user-oriented domain collection.',18,True,INK)
footer(s,10)

# 11
s=prs.slides.add_slide(blank); title(s,'10 • functional evidence','The feature set is broader than the original 12-slide deck')
items=[('Genre browse','Fantasy • Horror • Science fiction'),('Grouping','Author or publication year'),('Ordering','A–Z / Z–A'),('Advanced search','Author • year • genre • edition • languages • name type • title'),('Phase 2 views','Classification • names • titles • topics'),('User state','In-memory search list • reset • favorites persistence'),('Data quality','Missing values • multi-value fields • Unicode • leading-zero IDs'),('Reporting','Top 10 queries/results • >100 notification')]
for i,(h,b) in enumerate(items):
    col=i%2; row=i//2; x=0.75+col*6.15; y=1.95+row*1.03; box(s,x,y,5.75,0.88,h,b,fill=PALE,line=ACCENT,title_size=15,body_size=11)
footer(s,11)

# 12
s=prs.slides.add_slide(blank); title(s,'11 • acceptance','Real-data validation matches the required dataset and behavior')
views=[('records.csv','57,746'),('names.csv','117,873'),('titles.csv','77,280'),('topics.csv','77,919'),('classification.csv','57,844')]
for i,(f,n) in enumerate(views):
    box(s,0.75,1.85+i*0.73,4.7,0.58,f,n,fill=WHITE,line=GREEN,title_size=13,body_size=15)
box(s,6.0,1.85,6.1,3.2,'Behavior checks','Fantasy: 4,793\nHorror: 1,929\nScience Fiction: 9,356\n\nAuthor-role search: PASS\nUnicode title search: PASS\nMissing Publisher / ISBN: PASS\nMulti-value display: PASS\nA–Z / Z–A ordering: PASS',fill=PALE,line=ACCENT,title_size=20,body_size=14)
text(s,6.05,5.35,5.8,0.65,'Aggregation: 117,873 → 54,147 records\nThreshold: 101 inclusions → notification triggered',16,True,INK)
footer(s,12)

# 13
s=prs.slides.add_slide(blank); title(s,'12 • testing','Testing uses a layered evidence model')
# pyramid
levels=[('Unit','Focused domain + strategy assertions',0.9,1.75,3.0,0.72,ACCENT),('Integration','CSV/service boundaries + five official views',0.9,2.7,5.0,0.72,ACCENT2),('CLI / E2E','User-visible workflows + ordering',0.9,3.65,7.0,0.72,GOLD),('Real-data acceptance','Official dataset + assignment-specific criteria',0.9,4.6,9.2,0.72,GREEN)]
for h,b,x,y,w,hgt,c in levels: box(s,x,y,w,hgt,h,b,fill=WHITE,line=c,title_size=14,body_size=11)
text(s,10.5,2.0,2.0,2.8,'WHY LAYERS?\n\nFast tests find logic regressions.\n\nBroad tests catch wiring/schema problems.\n\nReal-data checks validate assumptions.',15,True,INK)
footer(s,13)

# 14
s=prs.slides.add_slide(blank); title(s,'13 • test results','33 / 33 automated tests pass')
text(s,0.85,2.0,4.0,0.6,'33 / 33',32,True,GREEN,PP_ALIGN.CENTER)
text(s,0.85,2.6,4.0,0.45,'PASSING TESTS',12,True,MUTED,PP_ALIGN.CENTER)
checks=[('Domain & model','PASS'),('Search strategies','PASS'),('Repositories','PASS'),('Aggregation','PASS'),('Phase 2 acceptance','PASS'),('Real five-view load','PASS'),('CLI / E2E','PASS'),('Threshold notification','PASS')]
for i,(h,r) in enumerate(checks):
    col=i%2; row=i//2; x=5.15+col*3.6; y=1.8+row*1.0; box(s,x,y,3.2,0.78,h,r,fill=WHITE,line=GREEN,title_size=13,body_size=12)
box(s,0.85,3.55,3.95,1.6,'Testing principle','Coverage is a diagnostic, not proof.\nAcceptance assertions tie evidence to actual assignment requirements.',fill=PALE,line=ACCENT,title_size=17,body_size=14)
footer(s,14)

# 15
s=prs.slides.add_slide(blank); title(s,'14 • algorithms','Simple algorithms, explicit complexity')
box(s,0.75,1.9,3.6,1.45,'LINEAR SEARCH','Filtering/search scans the in-memory sequence → O(n)',fill=WHITE,line=ACCENT,title_size=18,body_size=14)
box(s,4.85,1.9,3.6,1.45,'SORTING','Python stable sort → O(n log n)',fill=WHITE,line=GOLD,title_size=18,body_size=14)
box(s,8.95,1.9,3.6,1.45,'TOKEN COUNTS','Dictionary/Counter aggregation → O(n) expected',fill=WHITE,line=GREEN,title_size=18,body_size=14)
text(s,0.8,3.9,11.5,0.8,'Why not a database index? The brief explicitly asks for an in-memory educational application. The current linear approach is deterministic and easy to test.',19,True,INK)
box(s,0.8,5.05,11.5,0.9,'SCALING NOTE','For much larger or frequently changing collections, an indexed database or inverted search structure would become preferable.',fill=PALE,line=ACCENT,title_size=14,body_size=14)
footer(s,15)

# 16
s=prs.slides.add_slide(blank); title(s,'15 • quality','Clean coding improves correctness—not just readability')
cols=[('Naming','Explicit domain and service names make intent inspectable.'),('Small methods','Parsing, grouping, searching, and presentation have separate contracts.'),('Validation','Required ID/title checks happen before domain use.'),('Immutable state','Frozen dataclasses + tuples reduce accidental mutation.'),('Dependency injection','Repository abstractions make isolated tests practical.'),('Separation','CLI handles interaction; services handle business rules.')]
for i,(h,b) in enumerate(cols):
    col=i%3; row=i//3; x=0.75+col*4.05; y=1.95+row*1.75; box(s,x,y,3.7,1.45,h,b,fill=WHITE,line=ACCENT,title_size=16,body_size=13)
text(s,0.8,5.75,11.4,0.65,'Result: algorithms and data structures are easier to reason about because the code mirrors the processing pipeline.',18,True,INK)
footer(s,16)

# 17
s=prs.slides.add_slide(blank); title(s,'16 • discovery audit','Security discovery found no confirmed high-impact vulnerability')
box(s,0.75,1.9,3.65,2.1,'Confirmed reportable','None established during scoped discovery.',fill=WHITE,line=GREEN,title_size=20,body_size=18)
box(s,4.85,1.9,3.65,2.1,'Deferred proof gap','XML adapter reaches ElementTree.parse; no hostile-XML regression test or network trust boundary was established.',fill=WHITE,line=GOLD,title_size=20,body_size=15)
box(s,8.95,1.9,3.65,2.1,'Not applicable','No runtime SQL, shell, eval, SSRF, auth/tenant, unsafe deserialization, or dynamic-code sink found.',fill=WHITE,line=ACCENT,title_size=20,body_size=15)
text(s,0.8,4.55,11.5,0.8,'Important distinction: the XML item is a proof gap, not a confirmed vulnerability. Source code was not changed during the discovery pass.',17,True,INK)
box(s,0.8,5.55,11.5,0.8,'AUDIT ARTIFACTS','security-discovery/rank_input.jsonl • deep_review_input.jsonl • work_ledger.jsonl • repository_coverage_ledger.md • candidate/dedupe records',fill=PALE,line=ACCENT,title_size=13,body_size=12)
footer(s,17)

# 18
s=prs.slides.add_slide(blank); title(s,'17 • conclusion','The strongest result is traceability')
text(s,0.85,1.9,11.5,0.65,'Requirement → design decision → implementation → test → evidence',25,True,INK,PP_ALIGN.CENTER)
for i,(h,b,c) in enumerate([('Requirement','Real facet dataset',ACCENT),('Design','Aggregation + Repository',GOLD),('Implementation','Immutable domain + services',ACCENT),('Verification','33 tests + real acceptance',GREEN)]):
    x=0.8+i*3.05; box(s,x,3.0,2.65,1.6,h,b,fill=WHITE,line=c,title_size=17,body_size=14)
    if i<3: line(s,x+2.65,3.8,x+3.02,3.8)
text(s,0.85,5.25,11.4,0.65,'SOLID and patterns are valuable here because they isolate real sources of variation. They also carry an abstraction cost—and the evaluation should say so.',18,True,INK,PP_ALIGN.CENTER)
footer(s,18)

prs.save(OUT)
print(OUT)

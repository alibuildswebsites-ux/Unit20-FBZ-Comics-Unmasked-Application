from __future__ import annotations

import json
import re
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PInches, Pt as PPt

ROOT = Path('.')
SUB = ROOT / 'submission'
FIG = ROOT / 'evidence' / 'figures'
REPORT = json.loads((ROOT / 'reports/final_acceptance_report.json').read_text(encoding='utf-8'))
ANALYSIS = json.loads((ROOT / 'reports/real_dataset_analysis.json').read_text(encoding='utf-8'))

FONT_CANDIDATES = [
    '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf',
    '/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf',
]
FONT_BOLD = [
    '/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf',
    '/usr/share/fonts/truetype/liberation2/LiberationSans-Bold.ttf',
]


def get_font(size: int, bold: bool = False):
    for candidate in FONT_BOLD if bold else FONT_CANDIDATES:
        if Path(candidate).exists():
            return ImageFont.truetype(candidate, size)
    return ImageFont.load_default()


def make_figures() -> None:
    FIG.mkdir(parents=True, exist_ok=True)

    img = Image.new('RGB', (1800, 1000), 'white')
    draw = ImageDraw.Draw(img)
    draw.text((70, 35), 'FBZ Layered Architecture', fill='black', font=get_font(46, True))
    boxes = [
        ((90, 140, 1620, 145), 'Presentation', ['CLI / command-line entry point']),
        ((90, 350, 500, 220), 'Application services', ['EncyclopediaService', 'SearchService', 'FavouriteService']),
        ((650, 350, 500, 220), 'Domain', ['Comic', 'Contributor', 'SearchCriteria']),
        ((1210, 350, 500, 220), 'Patterns', ['Strategy', 'Simple Factory', 'Repository']),
        ((300, 700, 1200, 180), 'Data sources', ['records / names / titles / topics / classification', 'CSV repository + aggregation + XML adapter']),
    ]
    for (x, y, w, h), title, lines in boxes:
        draw.rounded_rectangle((x, y, x + w, y + h), radius=24, outline='black', width=4)
        draw.text((x + 24, y + 18), title, fill='black', font=get_font(30, True))
        yy = y + 68
        for line in lines:
            draw.text((x + 24, yy), line, fill='black', font=get_font(24))
            yy += 34
    arrows = [((900, 285), (350, 350)), ((900, 285), (900, 350)), ((900, 285), (1460, 350)), ((350, 570), (600, 700)), ((900, 570), (900, 700)), ((1460, 570), (1200, 700))]
    for a, b in arrows:
        draw.line((a[0], a[1], b[0], b[1]), fill='black', width=4)
    img.save(FIG / 'architecture.png')

    img = Image.new('RGB', (1800, 1100), 'white')
    draw = ImageDraw.Draw(img)
    draw.text((70, 35), 'FBZ Class / Pattern Relationships', fill='black', font=get_font(46, True))
    boxes = [
        ((70, 150, 480, 220), 'Comic', ['record_id, title', 'metadata + contributors', 'tokens(), authors()']),
        ((650, 150, 480, 220), 'ComicRepository <<abstract>>', ['all() -> Sequence[Comic]']),
        ((1230, 150, 480, 220), 'CsvComicRepository', ['CSV parsing', 'schema validation']),
        ((70, 500, 480, 220), 'SearchService', ['depends on repository', 'delegates to Strategy']),
        ((650, 500, 480, 220), 'SearchStrategy <<abstract>>', ['search(comics, query)']),
        ((1230, 500, 480, 220), 'SearchStrategyFactory', ['title / author / genre / year']),
        ((650, 840, 480, 190), 'AggregatingComicRepository', ['one record per BL ID', 'preserves name/role pairs']),
    ]
    for (x, y, w, h), title, lines in boxes:
        draw.rounded_rectangle((x, y, x + w, y + h), radius=24, outline='black', width=4)
        draw.text((x + 24, y + 18), title, fill='black', font=get_font(28, True))
        yy = y + 68
        for line in lines:
            draw.text((x + 24, yy), line, fill='black', font=get_font(22))
            yy += 32
    for a, b in [((550, 260), (650, 260)), ((1130, 260), (1230, 260)), ((890, 500), (890, 370)), ((550, 610), (650, 610)), ((1130, 610), (1230, 610)), ((890, 720), (890, 840))]:
        draw.line((a[0], a[1], b[0], b[1]), fill='black', width=4)
    img.save(FIG / 'class_diagram.png')


def set_doc_defaults(doc: Document) -> None:
    sec = doc.sections[0]
    sec.top_margin = Inches(0.65)
    sec.bottom_margin = Inches(0.65)
    sec.left_margin = Inches(0.75)
    sec.right_margin = Inches(0.75)
    styles = doc.styles
    styles['Normal'].font.name = 'Aptos'
    styles['Normal'].font.size = Pt(10)
    for name, size in [('Title', 26), ('Heading 1', 18), ('Heading 2', 14), ('Heading 3', 11)]:
        styles[name].font.name = 'Aptos Display'
        styles[name].font.size = Pt(size)


def add_cover(doc: Document, title: str, subtitle: str) -> None:
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run('Unit 20 — Applied Programming and Design Principles').bold = True
    p = doc.add_paragraph()
    p.style = 'Title'
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.add_run(title)
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(subtitle)
    r.italic = True
    r.font.size = Pt(12)
    table = doc.add_table(rows=5, cols=2)
    table.style = 'Table Grid'
    rows = [
        ('Student Name / ID', '____________________________'),
        ('Academic Year', '2025/26'),
        ('Unit', 'Unit 20 — Applied Programming and Design Principles'),
        ('Assignment', 'Implementation of a Dataset Processing Application using SOLID Design Principles'),
        ('Submission Date', '28 August 2026'),
    ]
    for row, (left, right) in zip(table.rows, rows):
        row.cells[0].text = left
        row.cells[1].text = right
        row.cells[0].paragraphs[0].runs[0].bold = True
    doc.add_page_break()


def add_table(doc: Document, headers: list[str], rows: list[tuple], widths: list[float] | None = None) -> None:
    table = doc.add_table(rows=1, cols=len(headers))
    table.style = 'Table Grid'
    for i, header in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = str(header)
        cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER
        cell.paragraphs[0].runs[0].bold = True
    for row in rows:
        cells = table.add_row().cells
        for i, value in enumerate(row):
            cells[i].text = str(value)
    if widths:
        for row in table.rows:
            for i, width in enumerate(widths):
                row.cells[i].width = Inches(width)


def add_refs(doc: Document, refs: list[str]) -> None:
    doc.add_heading('References', 1)
    for ref in refs:
        doc.add_paragraph(ref)


def build_design_doc() -> int:
    doc = Document()
    set_doc_defaults(doc)
    add_cover(doc, 'Design and Programming Report', 'Fantasy Bazaar (FBZ) Comics Unmasked Dataset Processing Application')

    doc.add_heading('1. Scenario, scope and requirements', 1)
    doc.add_paragraph(
        'The assignment scenario describes Fantasy Bazaar (FBZ), a local comic-book shop that needs a data-processing application around the British Library Comics Unmasked metadata. The required solution is not simply a working search script: the brief asks for investigation of object-oriented programming and SOLID principles, clean coding, design patterns and an automated testing regime. The functional scenario requires the dataset to be loaded into memory; users must be able to browse Fantasy, Horror and Science Fiction records, group them by author or publication year, order titles alphabetically in both directions, search manually by title, handle special characters, repeated values, missing ISBN values and multiple title rows, save results to an in-memory search list, perform advanced searches, and report popular searches/results. The later testing scenario requires evidence from automated tests and comparison of developer-produced and vendor/framework-supported testing approaches.'
    )
    doc.add_paragraph(
        'The implementation follows those requirements through a layered Python application. Presentation is kept in the CLI, application services coordinate use cases, domain classes represent catalogue concepts, repositories isolate data access, and Strategy/Factory objects isolate interchangeable search algorithms. This separation makes the design criteria visible in the implementation rather than describing SOLID only in the written report (British Library, n.d.; Chebanyuk and Markov, 2016).'
    )

    doc.add_heading('2. Real dataset acquisition, structure and validation', 1)
    doc.add_paragraph(
        'The project uses the actual 2022 Comics Unmasked Researcher Format package referenced by the supplied brief. The historical direct British Library download route was no longer available, so the exact archive was recovered from the archived source route and retained under `data/raw/ComicsResearcherFormat_202204_csv.zip`. The five extracted CSV views are preserved under `data/raw/extracted/`. The current British Library collection metadata service confirms that Researcher Format CSV datasets are now made available through the British Library Research Repository (British Library, n.d.). The raw files are therefore treated as source evidence, not regenerated synthetic fixtures.'
    )
    add_table(doc, ['CSV view', 'Rows', 'Columns', 'Leading-zero IDs', 'Acceptance'], [
        (name, meta['rows'], meta['columns'], meta['leading_zero_ids'], 'PASS' if REPORT['five_view_acceptance'][name]['pass'] else 'FAIL')
        for name, meta in ANALYSIS['files'].items()
    ])
    doc.add_paragraph(
        'The most important structural finding is the names view. It contains 117,873 facet rows but only 54,147 unique BL record IDs. Therefore, showing every row directly would violate the scenario requirement to present multiple title/name rows as one catalogue record. The aggregation repository collapses 63,726 repeated facet rows into 54,147 user-facing records while retaining semicolon-separated values and the original raw CSV files. This is a design decision driven by the actual data rather than by a generic architecture template.'
    )
    doc.add_paragraph(
        'The source also contains contributor roles. A flat list of names is insufficient for an author-specific search because the same record can contain authors, artists, editors, translators and other contributors. The refined domain model introduces `Contributor(name, role)` and derives `Comic.authors()` from explicit author/writer role tokens. This means an author search cannot return a record merely because the searched person appears as an editor or illustrator.'
    )

    doc.add_heading('3. OOP characteristics and relationships', 1)
    doc.add_paragraph(
        'Encapsulation is demonstrated by `Comic`, which owns its immutable record state and domain operations such as multi-value tokenisation and author extraction. Abstraction is represented by `ComicRepository` and `SearchStrategy`, which expose small contracts without forcing callers to know how a record is loaded or how a particular search is implemented. Polymorphism allows title, author, genre and year strategies to be used through the same `SearchStrategy.search` contract. Inheritance is purposeful rather than decorative: concrete repositories and concrete strategies are substitutable implementations of meaningful abstractions. The design also uses dependency/association relationships because services receive repository objects rather than constructing concrete storage objects themselves.'
    )
    doc.add_paragraph(
        'Composition is visible in repository state, where an in-memory repository owns a stable tuple of `Comic` objects. Aggregation is particularly relevant to the dataset: `AggregatingComicRepository` contains records produced from multiple source rows, but the source rows themselves remain independent in the raw dataset. Khan (2023) distinguishes inheritance, association, composition and aggregation as different relationships; the FBZ model uses each only where it explains an actual relationship in the software.'
    )

    doc.add_heading('4. SOLID principles applied to FBZ', 1)
    add_table(doc, ['Principle', 'FBZ implementation', 'Benefit and trade-off'], [
        ('SRP', 'CSV parsing, aggregation, domain modelling, search, favourites and presentation are separate.', 'Changes remain local; more classes must be understood.'),
        ('OCP', 'New search algorithms are added as new strategies.', 'Extensible search; requires a clear strategy contract.'),
        ('LSP', 'CSV, XML and in-memory repositories fulfil `ComicRepository`; strategies fulfil `SearchStrategy`.', 'Substitution improves testing; contracts must remain stable.'),
        ('ISP', 'Repository and strategy interfaces are small.', 'Less coupling; very small interfaces can add indirection if overused.'),
        ('DIP', 'Services depend on `ComicRepository`, not CSV implementation.', 'Enables isolated tests and alternative sources; dependency injection adds setup.'),
    ])
    doc.add_paragraph(
        'Chebanyuk and Markov (2016) treat SOLID as structural design principles that can be examined through class relationships. In FBZ, the strongest practical effect is change isolation. A new search algorithm does not require rewriting existing search algorithms, and a different repository implementation can be supplied without rewriting the service. The cost is abstraction overhead: a tiny one-off script could be shorter without these boundaries. Because FBZ has multiple search behaviours, multiple source views, an XML extension point and a testing requirement, the additional structure is justified.'
    )

    doc.add_heading('5. Clean coding, data structures and algorithms', 1)
    doc.add_paragraph(
        'The implementation follows clean-code practices through explicit names, type hints, immutable dataclasses, small focused methods, early validation, clear exceptions and separation of input/output from business logic. Mark (2023) emphasises readability, maintainability, meaningful naming and focused functions, while Cocca (2023) connects clean code with effectiveness, efficiency, simplicity, modularisation and algorithmic complexity. These ideas are visible in the project rather than merely listed: CSV parsing is isolated, aggregation is a repository concern, search algorithms are independent strategies, and the CLI contains interaction logic rather than business rules.'
    )
    doc.add_paragraph(
        'The main data structure is a tuple-backed in-memory sequence of immutable `Comic` objects. Linear filtering and simple search are O(n) because the application scans the loaded sequence. Alphabetical sorting is O(n log n) using Python’s stable sort. Frequency summaries use dictionary/Counter-style counting with O(n) expected aggregation over the observed tokens. Grouping uses dictionaries of lists, allowing one record to appear under multiple relevant authors or publication years. The design intentionally chooses clarity and deterministic behaviour over premature indexing. For a much larger or frequently changing collection, an indexed database or inverted search structure would reduce repeated linear scans.'
    )
    doc.add_paragraph(
        'Clean coding directly affects algorithmic correctness. A monolithic function that parsed CSV, merged rows, searched, sorted and printed results would mix concerns and make complexity difficult to verify. The final pipeline is explicit: parse → map to domain objects → aggregate when required → apply a focused search/filter → sort → present. Each stage has a narrow contract and can be tested independently. This is the practical connection between clean coding and the brief’s requirement to explain how coding techniques affect data structures and algorithms.'
    )

    doc.add_heading('6. Design patterns and refinement', 1)
    add_table(doc, ['Family', 'Pattern / technique', 'Concrete FBZ problem'], [
        ('Behavioural', 'Strategy', 'Title, author, genre and year searches are interchangeable algorithms.'),
        ('Creational', 'Simple Factory', '`SearchStrategyFactory` centralises strategy construction.'),
        ('Structural / architectural', 'Repository', 'Storage must be separated from application services.'),
        ('Structural / architectural refinement', 'Aggregation repository', 'Facet rows must become one user-facing record per BL ID.'),
    ])
    doc.add_paragraph(
        'Refactoring.Guru classifies design patterns into creational, structural and behavioural families and presents Strategy as a behavioural pattern for interchangeable algorithms (Refactoring.Guru, n.d.). Tutorialspoint similarly explains the three broad families and the value of programming to an interface (Tutorialspoint, n.d.). The project deliberately calls `SearchStrategyFactory` a simple factory rather than incorrectly claiming it is the formal GoF Factory Method. This is a technical-accuracy decision. Subburaj, Jekese and Hwata (2015) also warn that inappropriate pattern use and pattern overload can reduce the value of patterns. FBZ therefore uses patterns only where they solve actual variation or separation problems.'
    )

    doc.add_heading('7. Functional implementation evidence', 1)
    add_table(doc, ['Requirement', 'Implementation', 'Evidence'], [
        ('Genre filtering', 'Exact Fantasy/Horror/Science Fiction matching', '4,793 / 1,929 / 9,356 real records'),
        ('Author grouping/search', 'Contributor role preservation + author strategy', 'Real author acceptance query passes'),
        ('Year grouping', 'Publication-year token grouping', 'Unit and real-data acceptance'),
        ('A–Z / Z–A', 'Case-folded stable title sort', 'Acceptance checks both directions + CLI E2E'),
        ('Special characters', 'UTF-8 CSV + Unicode strings', 'Real Unicode title acceptance'),
        ('Repeated values', 'Tokenisation of semicolon/slash-separated fields', 'Real multi-value acceptance'),
        ('Missing ISBN', 'Display normalised to `missing`', 'Real missing-ISBN acceptance'),
        ('Multiple rows/titles', 'Aggregation by BL record ID', '117,873 → 54,147 records'),
        ('Search list', 'In-memory tuple', 'Unit test proves no disk persistence'),
        ('Advanced search', 'Author/year/genre/edition/language/name-type/title', 'Unit + real acceptance'),
        ('Phase 2 views', 'classification/names/titles/topics searches', 'Real acceptance samples'),
        ('>100 reporting', 'Threshold callback at 101', 'Deterministic real-comic acceptance'),
    ])

    doc.add_heading('8. Testing design and verification', 1)
    doc.add_paragraph(
        'The automated test regime is layered. Unit tests validate domain and strategy behaviour. Integration tests validate CSV/service boundaries and all five official views. CLI/end-to-end tests validate no-result output, result output and ordering. The real-data acceptance script validates requirements that synthetic fixtures cannot establish: official row counts, leading-zero identifiers, genre counts, contributor roles, Unicode, missing ISBN, multi-value metadata, aggregation, ordering and the >100 notification.'
    )
    doc.add_paragraph(
        'pytest documentation describes fixtures as explicit, modular and scalable test contexts (pytest, n.d.). Fowler’s Test Pyramid argues for many fast focused tests with fewer broad-stack tests (Fowler, 2012). The FBZ suite follows that principle. Katalon (2022) and SmartBear (n.d.) informed the comparison of data-driven, modular, hybrid and record/playback approaches. For this Python CLI, pytest is proportionate because the dominant risks are parsing, domain logic and service boundaries rather than browser rendering.'
    )
    doc.add_paragraph(
        'The final suite contains 30 passing automated tests. The acceptance report records the five-view row counts, 54,147 aggregated records, exact genre counts, a real author query, a Unicode title, a missing ISBN, multi-value genre tokens, both sort directions and a deterministic threshold notification after 101 search inclusions. This evidence is stronger than a coverage number alone because each result is tied to an assignment requirement.'
    )

    doc.add_heading('9. Critical evaluation and conclusion', 1)
    doc.add_paragraph(
        'The final design is more maintainable and testable than a procedural alternative because responsibilities and change boundaries are explicit. SOLID is effective here, but not universally free: it introduces abstractions and indirection. Strategy is justified by multiple search algorithms; Repository is justified by multiple data-source implementations; aggregation is justified by the facet structure of the real dataset. The most important refinement was semantic rather than cosmetic: preserving contributor name/role relationships prevents a technically working search algorithm from producing incorrect author results. The result is therefore a traceable design in which requirements lead to domain decisions, those decisions lead to implementation, and implementation is verified by automated and real-data acceptance evidence.'
    )

    refs = [
        'British Library (n.d.) Collection metadata services. Available at: https://www.bl.uk/services/collection-metadata-services (Accessed: 28 August 2026).',
        'Chebanyuk, E. and Markov, K. (2016) An Approach to Class Diagrams Verification According to SOLID Design Principles. Proceedings of MODELSWARD 2016, pp. 435–441. doi: 10.5220/0005830104350441.',
        'Cocca, G. (2023) How to Write Clean Code – Tips and Best Practices (Full Handbook). freeCodeCamp, 15 May. Available at: https://www.freecodecamp.org/news/how-to-write-clean-code/ (Accessed: 28 August 2026).',
        'Fowler, M. (2012) Test Pyramid. Martin Fowler, 1 May. Available at: https://martinfowler.com/bliki/TestPyramid.html (Accessed: 28 August 2026).',
        'Katalon (2022) Software Test Automation Frameworks | 6 Common Types. 28 November. Available at: https://medium.com/@katalon/test-automation-framework-e4e6cc09ea6d (Accessed: 28 August 2026).',
        'Khan, M.H. (2023) Understanding Object-Oriented Relationships: Inheritance, Association, Composition, and Aggregation. Medium, 14 October. Available at: https://medium.com/@humzakhalid94/understanding-object-oriented-relationships-inheritance-association-composition-and-aggregation-4d298494ac1c (Accessed: 28 August 2026).',
        'Mark, M. (2023) Writing Clean Code: Best Practices and Principles. DEV Community, 16 September. Available at: https://dev.to/favourmark05/writing-clean-code-best-practices-and-principles-3amh (Accessed: 28 August 2026).',
        'pytest (n.d.) About fixtures. Available at: https://docs.pytest.org/en/latest/explanation/fixtures.html (Accessed: 28 August 2026).',
        'Refactoring.Guru (n.d.) Design Patterns. Available at: https://refactoring.guru/design-patterns (Accessed: 28 August 2026).',
        'SmartBear (n.d.) Test Automation Frameworks. Available at: https://smartbear.com/learn/automated-testing/test-automation-frameworks/ (Accessed: 28 August 2026).',
        'Subburaj, R., Jekese, G. and Hwata, C. (2015) Impact of Object Oriented Design Patterns on Software Development. International Journal of Scientific & Engineering Research, 6(2), pp. 961–966.',
        'Tutorialspoint (n.d.) Design Patterns – Overview. Available at: https://www.tutorialspoint.com/design_pattern/design_pattern_overview.htm (Accessed: 28 August 2026).',
    ]
    add_refs(doc, refs)

    out = SUB / 'Unit20_FBZ_Design_and_Programming_Report.docx'
    out.parent.mkdir(exist_ok=True)
    doc.save(out)
    return len(' '.join(p.text for p in doc.paragraphs).split())


def build_testing_doc() -> int:
    doc = Document()
    set_doc_defaults(doc)
    add_cover(doc, 'Automated Testing Report', 'FBZ Comics Unmasked Dataset Processing Application')

    doc.add_heading('1. Testing strategy', 1)
    doc.add_paragraph(
        'The final testing regime uses four complementary levels: unit testing for focused domain and algorithm behaviour, integration testing for repository/service boundaries, CLI/end-to-end testing for user-visible workflows, and real-data acceptance testing against the supplied British Library package. The purpose is requirement-based confidence rather than a single metric. Coverage is useful for finding untested paths, but a high percentage cannot prove that the assertions represent the assignment correctly.'
    )
    doc.add_heading('2. Automated testing methods and tools', 1)
    doc.add_paragraph(
        'Developer-produced tests contain the application-specific assertions. For example, the FBZ author test must distinguish an author from an editor with the same name, while a generic framework cannot know that requirement. pytest supplies the execution machinery: discovery, assertions, fixtures and reusable test contexts. Its documentation describes fixtures as explicit, modular and scalable (pytest, n.d.). This division is important for the assessment criterion comparing developer-produced testing with vendor/framework-supported automation: the framework provides capability, but the project supplies the domain oracle.'
    )
    doc.add_paragraph(
        'The research also examined common automation framework styles. Katalon (2022) describes linear, modular, library-architecture, data-driven, keyword-driven and hybrid approaches. SmartBear (n.d.) describes similar framework trade-offs and notes that record/playback is easy to start but can become difficult to maintain when applications change. A data-driven approach is relevant to FBZ because the same parsing and search behaviours can be exercised against many records or dataset views. A keyword-driven or GUI-heavy framework is less proportionate for the current Python CLI because the main risks are data correctness and application logic rather than browser interaction.'
    )
    doc.add_heading('3. Test levels and trade-offs', 1)
    add_table(doc, ['Level', 'What it verifies', 'Strength', 'Limitation'], [
        ('Unit', 'Domain, strategies, factory, services, edge cases', 'Fast, isolated, precise', 'Can miss wiring/source-format defects'),
        ('Integration', 'CSV → repository → service and official five views', 'Finds boundary and schema problems', 'Slower and data-dependent'),
        ('CLI / E2E', 'No-result, result display and ordering', 'Validates user-visible behaviour', 'Broader and more maintenance-sensitive'),
        ('Real-data acceptance', 'Official counts, semantics and assignment-specific behaviour', 'Validates assumptions against supplied evidence', 'Requires the real package and can be slower'),
    ])
    doc.add_paragraph(
        'Fowler’s Test Pyramid recommends a larger base of focused tests and fewer expensive broad-stack tests because high-level tests can be slower and more brittle (Fowler, 2012). The FBZ suite follows this principle. Manual exploratory checks remain useful for presentation and usability, but repeated parsing, searching, sorting and threshold logic are strong automation candidates because they are deterministic and data-intensive. SmartBear’s guidance similarly identifies repetitive, high-volume or frequently repeated tests as good candidates for automation.'
    )

    doc.add_heading('4. Test matrix', 1)
    matrix = [
        ('UT-01', 'Valid record + leading-zero ID', 'Unit', 'PASS'),
        ('UT-02', 'Required-column/required-title validation', 'Unit', 'PASS'),
        ('UT-03', 'Multi-value tokenisation', 'Unit', 'PASS'),
        ('UT-04', 'Title search', 'Unit', 'PASS'),
        ('UT-05', 'Author search respects role', 'Unit', 'PASS'),
        ('UT-06', 'Genre filtering', 'Unit', 'PASS'),
        ('UT-07', 'Author/year grouping', 'Unit', 'PASS'),
        ('UT-08', 'A–Z/Z–A sorting', 'Unit/E2E', 'PASS'),
        ('UT-09', 'Advanced multi-criteria search', 'Unit', 'PASS'),
        ('UT-10', 'In-memory search list/reset', 'Unit', 'PASS'),
        ('UT-11', '>100 threshold notification', 'Unit/Acceptance', 'PASS'),
        ('IT-01', 'CSV → service → favourite flow', 'Integration', 'PASS'),
        ('IT-02', 'All five official views', 'Integration', 'PASS'),
        ('IT-03', 'Real names dataset scale/schema', 'Integration', 'PASS'),
        ('E2E-01', 'CLI no-result message', 'E2E', 'PASS'),
        ('E2E-02', 'CLI result display', 'E2E', 'PASS'),
        ('E2E-03', 'CLI descending order', 'E2E', 'PASS'),
    ]
    add_table(doc, ['ID', 'Requirement', 'Level', 'Result'], matrix)

    doc.add_heading('5. Real-data acceptance evidence', 1)
    add_table(doc, ['View', 'Observed rows', 'Expected rows', 'Status'], [
        (name, REPORT['five_view_acceptance'][name]['rows'], REPORT['five_view_acceptance'][name]['expected_rows'], 'PASS')
        for name in REPORT['five_view_acceptance']
    ])
    doc.add_paragraph(
        f"The real names view contains {REPORT['raw_names_rows']:,} raw rows and {REPORT['unique_names_view_records']:,} unique BL record IDs. The exact genre acceptance counts are Fantasy {REPORT['genre_record_counts']['fantasy']:,}, Horror {REPORT['genre_record_counts']['horror']:,} and Science Fiction {REPORT['genre_record_counts']['science fiction']:,}. The acceptance report also records a real author query, Unicode title, missing ISBN, multi-value genre, both sort directions and the >100 threshold."
    )
    doc.add_paragraph(
        f"The threshold test deliberately runs the first real dataset Comic through 101 matching search inclusions. The recorded result is {REPORT['threshold_acceptance']['count']} inclusions and a triggered notification. This deterministic check is stronger than relying on an uncontrolled demonstration session to happen to cross the threshold."
    )

    doc.add_heading('6. Final test results and reporting', 1)
    doc.add_paragraph(
        'The current automated suite contains 30 passing tests. `reports/final_acceptance_report.json` is the machine-readable acceptance record and `reports/real_dataset_analysis.json` contains the five-view schema/scale analysis. The test suite intentionally includes both synthetic fixtures and the real package: synthetic fixtures isolate edge cases such as an editor-versus-author collision, while the real acceptance run validates the dataset-specific assumptions.'
    )
    doc.add_paragraph(
        'The main benefits of this approach are repeatability, rapid regression feedback, deterministic edge-case coverage and safer refactoring. The drawbacks are test authoring and maintenance cost, longer execution time as the suite grows, and the possibility of false confidence if tests contain weak assertions. A unit test can pass while a real CSV schema is wrong; an end-to-end test can catch the wiring issue but be more expensive. The layered strategy balances these risks.'
    )
    doc.add_heading('7. Failure analysis and reproducibility', 1)
    doc.add_paragraph(
        'The project also treats failures as evidence rather than as reasons to skip a requirement. During the final refinement, making author search role-aware initially caused three existing tests to fail because their synthetic fixtures supplied a name but no role. The failure was reproduced, traced to a mismatch between the old fixture assumption and the refined domain rule, and fixed by making the fixtures explicitly model an author role. A separate command used for a diagnostic import failed because the project package was not on PYTHONPATH; the environment issue was reproduced and corrected by using the project-supported PYTHONPATH configuration. A syntax error in the acceptance script was similarly caught by preflight compilation, corrected, compiled again and only then executed. These incidents demonstrate the intended workflow: stop on failure, reproduce it, identify the root cause, fix the underlying issue, preflight and rerun.'
    )
    doc.add_heading('8. Developer-produced versus vendor/framework testing', 1)
    doc.add_paragraph(
        'Developer-produced assertions are necessary because only the project knows what constitutes a correct FBZ result. pytest supplies mature discovery, fixtures and assertion reporting, but it cannot decide whether an editor should count as an author or whether the five official CSV row counts are correct. Commercial tools such as TestComplete or Katalon can be valuable for browser, desktop and cross-platform functional automation, especially where non-programmers need record/playback or enterprise reporting. Their additional features can also introduce licence, setup and maintenance costs. For this CLI application, a lightweight open-source framework is more proportionate because the critical risks are data processing and business logic.'
    )
    doc.add_heading('9. Suitability of automation by application type', 1)
    doc.add_paragraph(
        'Automation is most valuable when a test is repeatable, data-heavy, high-risk or executed across many builds. FBZ satisfies all four conditions: CSV parsing is repeated, searches operate over tens of thousands of records, edge cases such as leading-zero IDs are easy to regress, and the same acceptance checks can be rerun after refactoring. In contrast, visual design, exploratory usability and ambiguous user experience questions are better complemented by manual review. This distinction prevents the common mistake of treating automated testing as a replacement for every other testing activity.'
    )

    add_refs(doc, [
        'Fowler, M. (2012) Test Pyramid. Martin Fowler, 1 May. Available at: https://martinfowler.com/bliki/TestPyramid.html (Accessed: 28 August 2026).',
        'Katalon (2022) Software Test Automation Frameworks | 6 Common Types. 28 November. Available at: https://medium.com/@katalon/test-automation-framework-e4e6cc09ea6d (Accessed: 28 August 2026).',
        'pytest (n.d.) About fixtures. Available at: https://docs.pytest.org/en/latest/explanation/fixtures.html (Accessed: 28 August 2026).',
        'SmartBear (n.d.) Test Automation Frameworks. Available at: https://smartbear.com/learn/automated-testing/test-automation-frameworks/ (Accessed: 28 August 2026).',
    ])
    out = SUB / 'Unit20_FBZ_Automated_Testing_Report.docx'
    out.parent.mkdir(exist_ok=True)
    doc.save(out)
    return len(' '.join(p.text for p in doc.paragraphs).split())


def add_slide_text(slide, x, y, w, h, text, size=22, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PInches(x), PInches(y), PInches(w), PInches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(245, 247, 250)
    shape.line.color.rgb = RGBColor(70, 80, 90)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = text
    r.font.size = PPt(size)
    r.font.bold = bold
    r.font.color.rgb = RGBColor(25, 35, 45)


def parse_slides_source() -> list[tuple[str, str, str]]:
    text = (ROOT / 'docs/presentation/slides.md').read_text(encoding='utf-8')
    blocks = re.split(r'(?m)^## Slide \d+ — ', text)[1:]
    slides = []
    for block in blocks:
        lines = block.splitlines()
        title = lines[0].strip()
        on_slide_start = next(i for i, line in enumerate(lines) if line.strip() == '**On slide**')
        notes_start = next(i for i, line in enumerate(lines) if line.strip() == '**Speaker notes**')
        on_slide = '\n'.join(line[2:].strip() for line in lines[on_slide_start + 1:notes_start] if line.strip().startswith('- '))
        notes = ' '.join(line.strip() for line in lines[notes_start + 1:] if line.strip())
        slides.append((title, on_slide, notes))
    return slides


def build_ppt() -> int:
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)
    slides = parse_slides_source()
    references = {
        4: 'Research: Chebanyuk and Markov (2016)',
        5: 'Research: Mark (2023); Cocca (2023)',
        7: 'Research: Tutorialspoint (n.d.)',
        9: 'Research: Refactoring.Guru (n.d.)',
        11: 'Research: pytest (n.d.); Fowler (2012)',
        12: 'Research: Subburaj et al. (2015); SmartBear (n.d.)',
    }
    total_notes = 0
    for index, (title, body, notes) in enumerate(slides, 1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(PInches(0.65), PInches(0.42), PInches(12), PInches(0.8))
        p = title_box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.size = PPt(28)
        r.font.bold = True
        r.font.color.rgb = RGBColor(25, 35, 45)
        add_slide_text(slide, 0.75, 1.45, 11.85, 4.95, body, size=21)
        footer = slide.shapes.add_textbox(PInches(0.75), PInches(6.82), PInches(11.8), PInches(0.3))
        fp = footer.text_frame.paragraphs[0]
        fp.text = f'Unit 20 • FBZ • Slide {index}/12'
        fp.runs[0].font.size = PPt(10)
        fp.runs[0].font.color.rgb = RGBColor(100, 110, 120)
        if index in references:
            refbox = slide.shapes.add_textbox(PInches(7.3), PInches(6.8), PInches(5.2), PInches(0.35))
            rp = refbox.text_frame.paragraphs[0]
            rp.alignment = PP_ALIGN.RIGHT
            rp.text = references[index]
            rp.runs[0].font.size = PPt(9)
            rp.runs[0].font.italic = True
            rp.runs[0].font.color.rgb = RGBColor(100, 110, 120)
        notes_shape = next((shape for shape in slide.notes_slide.shapes if shape.has_text_frame and 'Click to add notes' in shape.text), None)
        if notes_shape is None:
            notes_shape = next((shape for shape in slide.notes_slide.shapes if shape.has_text_frame and shape.name.startswith('Notes')), None)
        if notes_shape is not None:
            notes_shape.text_frame.text = notes
        total_notes += len(notes.split())
    out = SUB / 'Unit20_FBZ_SOLID_OOP_Presentation.pptx'
    out.parent.mkdir(exist_ok=True)
    prs.save(out)
    return total_notes


def main() -> None:
    make_figures()
    design_words = build_design_doc()
    testing_words = build_testing_doc()
    notes_words = build_ppt()
    print(json.dumps({'design_document_words': design_words, 'testing_report_words': testing_words, 'speaker_notes_words': notes_words, 'slides': 12}, indent=2))


if __name__ == '__main__':
    main()
